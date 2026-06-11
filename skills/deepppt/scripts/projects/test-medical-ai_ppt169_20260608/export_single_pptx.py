#!/usr/bin/env python
"""Render SVG to PNG using only stdlib + Pillow (text-based approximation), 
then embed into PPTX. For real SVG rendering, we convert via a web browser approach."""
import sys, os
import subprocess, base64, io

sys.path = [p for p in sys.path if '.workbuddy' not in p and 'workbuddy' not in p.lower()]

from PIL import Image, ImageDraw, ImageFont, ImageColor
from pptx import Presentation
from pptx.util import Emu

PROJECT = os.path.dirname(os.path.abspath(__file__))
SVG_FILE = os.path.join(PROJECT, 'svg_output', '01_ai_medical_scenarios.svg')
OUTPUT = os.path.join(PROJECT, 'exports', 'test_medical_ai.pptx')
png_path = os.path.join(PROJECT, 'svg_output', '_slide.png')

os.makedirs(os.path.join(PROJECT, 'exports'), exist_ok=True)

# Better approach: use the image_extract SVG which is already a well-formed image.
# Since python-pptx can't embed SVG directly, and cairosvg needs Cairo, 
# let's create a proper PNG representation.

# First, read the SVG to check it exists
with open(SVG_FILE, 'r', encoding='utf-8') as f:
    svg = f.read()

print(f'SVG file size: {len(svg)} bytes')

# Create a white 1280x720 base image and draw a stylized version
img = Image.new('RGB', (2560, 1440), (248, 249, 252))
draw = ImageDraw.Draw(img)

# Find a font
font_path = None
for candidate in [
    r'C:\Windows\Fonts\msyh.ttc',   # Microsoft YaHei
    r'C:\Windows\Fonts\msyhbd.ttc', # Microsoft YaHei Bold
    r'C:\Windows\Fonts\arial.ttf',
]:
    if os.path.exists(candidate):
        font_path = candidate
        break

if font_path:
    print(f'Using font: {font_path}')
    font_large = ImageFont.truetype(font_path, 48)
    font_title = ImageFont.truetype(font_path, 20)
    font_body = ImageFont.truetype(font_path, 24)
    font_small = ImageFont.truetype(font_path, 18)
    font_icon = ImageFont.truetype(font_path, 36)
else:
    font_large = ImageFont.load_default()
    font_title = ImageFont.load_default()
    font_body = ImageFont.load_default()
    font_small = ImageFont.load_default()
    font_icon = ImageFont.load_default()

print(f'Font size check - large: {font_large.size}')

# Top red accent bar
draw.rectangle([0, 0, 2560, 8], fill=(233, 0, 47))

# Title area
draw.rectangle([112, 80, 2456, 216], fill=(255, 255, 255))
draw.rectangle([112, 80, 6+112, 136], fill=(233, 0, 47))  # left accent bar

# Title text
draw.text((164, 100), "Medical AI Applications", fill=(136, 136, 136), font=font_title)
draw.text((164, 140), "智慧医疗 · AI典型应用场景配图", fill=(29, 29, 26), font=font_large)

# Card configurations: 6 cards, each with x,y,w,h
cards = [
    # (x, y, w, h, number, title, color, bg_color, items)
]

colors = [
    (233, 0, 47),    # red
    (37, 99, 235),   # blue
    (5, 150, 105),   # green
    (217, 119, 6),   # amber
    (124, 58, 237),  # purple
    (8, 145, 178),   # cyan
]

bg_colors = [
    (255, 245, 247),
    (239, 246, 255),
    (240, 253, 244),
    (255, 251, 235),
    (243, 232, 255),
    (236, 254, 255),
]

titles = [
    "AI智能导诊分诊", "AI辅助报告生成", "术中智能引导",
    "用药安全智能审核", "全病程随访管理", "AI医疗数据中心"
]

card_w = 744
card_h = 520
row1_y = 264
row2_y = 832
gap = 24

card_positions = [
    (112, row1_y), (868, row1_y), (1624, row1_y),
    (112, row2_y), (868, row2_y), (1624, row2_y),
]

# Draw 6 cards
for i, (cx, cy) in enumerate(card_positions):
    color = colors[i]
    bg = bg_colors[i]
    
    # Card background
    draw.rounded_rectangle([cx, cy, cx+card_w, cy+card_h], radius=24, fill=(255, 255, 255))
    
    # Accent circle with number
    circle_cx = cx + 50
    circle_cy = cy + 40
    draw.ellipse([circle_cx-24, circle_cy-24, circle_cx+24, circle_cy+24], fill=color)
    draw.text((circle_cx-8, circle_cy-12), str(i+1), fill=(255, 255, 255), font=font_icon)
    
    # Title
    draw.text((circle_cx+50, circle_cy-10), titles[i], fill=(29, 29, 26), font=font_body)
    
    # Subtitle
    subtitles = [
        "NLP驱动·精准分诊",
        "多模态·结构化报告",
        "实时引导·风险预警",
        "全链路·用药风控",
        "AI随访·全程管理",
        "数据汇聚·训推一体"
    ]
    draw.text((circle_cx+50, circle_cy+24), subtitles[i], fill=(150, 150, 150), font=font_small)
    
    # Process flow items (3 steps)
    step_y = cy + 100
    steps = [
        [("AI意图识别", "NLP引擎"), ("知识图谱匹配", "科室映射"), ("精准推荐挂号", "分诊引擎")],
        [("影像报告", "CT/MRI结构化"), ("病理报告", "组织学解析"), ("AI审核", "诊断建议")],
        [("🔬 内镜视野", "实时病灶识别"), ("⚡ 预警系统", "风险预警"), ("导航系统", "路径规划")],
        [("开方审核", "禁忌/相互作用"), ("剂量校验", "剂量/频次"), ("药师调配", "全链路风控")],
        [("出院跟踪", "AI随访计划"), ("7天康复", "康复跟踪"), ("30天复查", "复查提醒")],
        [("数据汇聚", "多源融合"), ("模型训练", "医疗大模型"), ("数据安全", "隐私计算")],
    ]
    
    step_colors = [color, (color[0]^50, color[1]^50, min(color[2]+50, 255)), (color[0]^100, color[1]^100, color[2]^100)]
    
    for j, (step_title, step_desc) in enumerate(steps[i]):
        sx = cx + 32 + j * 240
        draw.rounded_rectangle([sx, step_y, sx+224, step_y+100], radius=12, fill=(255,255,255), outline=color, width=2)
        draw.text((sx+112, step_y+20), step_title, fill=color, font=font_small, anchor="mt")
        draw.text((sx+112, step_y+55), step_desc, fill=(100, 100, 100), font=ImageFont.truetype(font_path, 16) if font_path else font_small, anchor="mt")
        
        # Arrow between steps
        if j < 2:
            ax = sx + 224 + 4
            draw.text((ax+4, step_y+35), "→", fill=(200,200,200), font=font_body)
    
    # Bottom info area
    info_y = cy + 230
    draw.rounded_rectangle([cx+16, info_y, cx+card_w-16, cy+card_h-16], radius=16, fill=bg)
    
    info_texts = [
        [
            "患者描述症状 → NLP意图识别 → 知识图谱映射",
            "精准推送对应科室专家、匹配门诊时段",
            "减少挂错号率 60%，缩短候诊时间 40%"
        ],
        [
            "图像智能识别→结构化描述→初步诊断建议",
            "支持 CT/MRI/PET/超声/病理等多模态数据",
            "报告生成效率提升 70%，医生审核后即可发布"
        ],
        [
            "术中实时病灶识别与边界勾画",
            "关键结构（血管/神经）实时标注与预警",
            "手术路径规划优化，并发症减少 35%"
        ],
        [
            "历史病历+处方方案 → 药物相互作用智能检测",
            "用药剂量 / 频次 / 过敏史自动校验",
            "多节点风控：开方→审核→调配→用药全链路",
            "用药差错减少 55%，门诊药房审核效率提升 3 倍"
        ],
        [
            "AI制定个性化随访计划，自动推送注意事项",
            "语音/图文随访，智能识别患者异常反馈",
            "异常指标自动预警，触发复诊/紧急干预",
            "随访管理率提升 80%，再入院率降低 25%"
        ],
        [
            "影像/病历/检验/基因/可穿戴设备多源汇聚",
            "基于昇腾AI基础软硬件，构建医疗领域大模型",
            "数据安全合规：联邦学习 + 隐私计算",
            "支持院内/区域/云三种部署模式"
        ],
    ]
    
    dot_color = color
    li_y = info_y + 20
    for k, text in enumerate(info_texts[i]):
        li = int(li_y)
        draw.ellipse([cx+36, li_y+6, cx+46, li_y+16], fill=dot_color)
        draw.text((cx+56, li_y), text, fill=(85, 85, 85), font=ImageFont.truetype(font_path, 20) if font_path else font_small)
        li_y += 42

# Footer
draw.rectangle([112, 1380, 2456, 1428], fill=(245, 245, 250))
draw.text((128, 1390), "华为LightAI · 智慧医疗解决方案 | 基于昇腾AI基础软硬件平台 | 覆盖诊前-诊中-诊后全流程", fill=(180, 180, 180), font=font_small)

img.save(png_path, "PNG")
print(f'Rendered PNG saved to: {png_path}')
print(f'Image size: {img.size}')

# Create PPTX with the rendered image
prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(png_path, Emu(0), Emu(0), Emu(12192000), Emu(6858000))

prs.save(OUTPUT)
os.remove(png_path)

print(f'PPTX saved to: {OUTPUT}')
print(f'File size: {os.path.getsize(OUTPUT) / 1024:.1f} KB')
