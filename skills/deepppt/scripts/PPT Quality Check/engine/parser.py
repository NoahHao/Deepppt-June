"""
PPTX 文本提取器 — 从 PPTX 文件中提取每页的标题、正文、全部文本。

支持两种模式：
  - extract_title_and_body(): 提取标题+正文（兼容 gate_checker）
  - extract_all_texts(): 提取每页全部文本块（供 content_quality_checker 使用）
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
import re


def _get_shape_texts(shape):
    """提取 shape 中所有文本段落。返回文本列表和字号列表。"""
    texts = []
    font_sizes = []
    if not shape.has_text_frame:
        return texts, font_sizes
    for para in shape.text_frame.paragraphs:
        full_text = para.text.strip()
        if full_text:
            texts.append(full_text)
        for run in para.runs:
            if run.font.size:
                font_sizes.append(run.font.size)
    return texts, font_sizes


def extract_title_and_body(pptx_path):
    """
    从 PPTX 中提取每页的标题和正文。
    标题检测策略（多级降级）：
      1. 占位符类型为 TITLE → 直接作为标题
      2. 没有 TITLE 占位符 → 取全页最大字号的文本块作为标题
      3. 仍无 → 标题为空，正文保留全部文本
    """
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []
        is_title_from_placeholder = False

        # 收集所有 shape 信息
        shape_info = []
        for shape in slide.shapes:
            texts, font_sizes = _get_shape_texts(shape)
            if not texts:
                continue

            is_title_placeholder = False
            if shape.is_placeholder:
                # 从 XML 读取 placeholder type: 'title', 'ctrTitle', 'body' 等
                ph_elem = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}ph')
                if ph_elem is not None:
                    ph_type = ph_elem.get('type', '')
                    is_title_placeholder = ph_type in ('title', 'ctrTitle')

            max_fs = max(font_sizes) if font_sizes else Pt(0)
            combined = ' '.join(texts)

            shape_info.append({
                'text': combined,
                'texts': texts,
                'max_font_size': max_fs,
                'is_title_placeholder': is_title_placeholder,
            })

        # 策略 1: 用 TITLE 占位符的内容作为标题
        for info in shape_info:
            if info['is_title_placeholder'] and info['text']:
                title = info['text']
                is_title_from_placeholder = True
                break

        # 策略 2: 没有 TITLE 占位符，取最大字号
        if not title:
            max_fs = Pt(0)
            for info in shape_info:
                if info['max_font_size'] > max_fs and info['text']:
                    max_fs = info['max_font_size']
                    title = info['text']

        # 构建正文（排除已作为标题的文本）
        for info in shape_info:
            text = info['text']
            if is_title_from_placeholder and info['is_title_placeholder']:
                continue  # 跳过标题占位符
            if text and text != title:
                body_parts.append(text)

        # 清理标题
        title = re.sub(r'^\d+\s*[\.\-\|]\s*', '', title).strip()
        body = '\n'.join(body_parts).strip()

        slides_data.append({
            'page': i,
            'title': title,
            'body': body,
        })

    return slides_data


def extract_all_texts(pptx_path):
    """
    从 PPTX 中提取每页的全部文本结构。
    返回: [{'page': int, 'shapes': [{'text': str, 'font_size_pt': float, 'left': int, 'top': int}, ...]}, ...]
    """
    prs = Presentation(pptx_path)
    all_data = []

    for i, slide in enumerate(prs.slides, start=1):
        shapes_data = []
        for shape in slide.shapes:
            texts, font_sizes = _get_shape_texts(shape)
            if not texts:
                continue

            full_text = ' '.join(texts)
            max_fs = max(font_sizes) if font_sizes else Pt(0)
            # 安全地将 Emu 转换为浮点数 pt
            if hasattr(max_fs, 'pt'):
                fs_pt = max_fs.pt
            elif isinstance(max_fs, (int, float)):
                fs_pt = max_fs / 12700  # EMU → pt
            else:
                fs_pt = 0

            shapes_data.append({
                'text': full_text,
                'font_size_pt': round(fs_pt, 1),
                'left_emu': shape.left,
                'top_emu': shape.top,
            })

        all_data.append({
            'page': i,
            'shapes': shapes_data,
        })

    return all_data
