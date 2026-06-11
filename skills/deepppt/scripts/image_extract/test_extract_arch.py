#!/usr/bin/env python3
"""
test_extract_arch.py -- Extract architecture diagrams from PPTX slide 1
and create a new PPT with white background.
Usage:
  python test_extract_arch.py <pptx_path> [output_pptx]
"""

import sys
import re
import os
import tempfile
import zipfile
from pathlib import Path
from io import BytesIO

from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor


def parse_slide1_images(pptx_path):
    """Parse slide1 to find all image shapes with position/size in EMU."""
    with zipfile.ZipFile(pptx_path, 'r') as z:
        slide_xml = z.read('ppt/slides/slide1.xml').decode('utf-8', errors='replace')
        try:
            rels_xml = z.read('ppt/slides/_rels/slide1.xml.rels').decode('utf-8', errors='replace')
        except Exception:
            rels_xml = ""

    shapes = []
    tag_starts = list(re.finditer(r'<(p:sp|p:pic|p:graphicFrame)\b', slide_xml))

    for start_m in tag_starts:
        tag_name = start_m.group(1)
        start_pos = start_m.start()

        depth = 1
        search_pos = start_m.end()
        end_pos = -1
        open_tag = f'<{tag_name}'
        close_tag = f'</{tag_name}>'

        while depth > 0 and search_pos < len(slide_xml):
            next_open = slide_xml.find(open_tag, search_pos)
            next_close = slide_xml.find(close_tag, search_pos)
            if next_close == -1:
                break
            if next_open != -1 and next_open < next_close:
                depth += 1
                search_pos = next_open + len(open_tag)
            else:
                depth -= 1
                search_pos = next_close + len(close_tag)
                if depth == 0:
                    end_pos = next_close + len(close_tag)

        if end_pos == -1:
            continue

        block = slide_xml[start_pos:end_pos]

        blip_match = re.search(r'<a:blip[^>]*r:embed="([^"]+)"', block)
        if not blip_match:
            continue

        embed_id = blip_match.group(1)

        target_match = re.search(
            rf'Id="{re.escape(embed_id)}"[^>]*Target="([^"]+)"',
            rels_xml
        )
        if not target_match:
            continue

        target = target_match.group(1)
        internal_name = Path(target).name

        spPr_match = re.search(r'<p:spPr\b[^>]*>.*?</p:spPr>', block, re.DOTALL)
        if not spPr_match:
            spPr_match = re.search(r'<p:grpSpPr\b[^>]*>.*?</p:grpSpPr>', block, re.DOTALL)
        if not spPr_match:
            continue

        spPr = spPr_match.group()
        xfrm_match = re.search(r'<a:xfrm\b.*?</a:xfrm>', spPr, re.DOTALL)
        if not xfrm_match:
            continue

        xfrm_text = xfrm_match.group()
        off_match = re.search(r'<a:off[^>]*x="(\d+)"[^>]*y="(\d+)"', xfrm_text)
        ext_match = re.search(r'<a:ext[^>]*cx="(\d+)"[^>]*cy="(\d+)"', xfrm_text)

        x = int(off_match.group(1)) if off_match else 0
        y = int(off_match.group(2)) if off_match else 0
        cx = int(ext_match.group(1)) if ext_match else 0
        cy = int(ext_match.group(2)) if ext_match else 0

        slide_w = 12192000
        slide_h = 6858000
        coverage = (cx * cy) / (slide_w * slide_h)
        is_bg = coverage >= 0.6

        shapes.append({
            'internal_name': internal_name,
            'x': x, 'y': y,
            'cx': cx, 'cy': cy,
            'coverage': coverage,
            'is_background': is_bg,
        })

    return shapes


def process_image_for_white_bg(img_data):
    """Process image data to composite on white background if needed."""
    img = Image.open(BytesIO(img_data))

    if img.mode in ('RGBA', 'LA', 'PA'):
        bg = Image.new('RGB', img.size, (255, 255, 255))
        if img.mode == 'RGBA':
            bg.paste(img, mask=img.split()[3])
        elif img.mode == 'PA':
            img = img.convert('RGBA')
            bg.paste(img, mask=img.split()[3])
        else:
            bg.paste(img)
        tmp_file = Path(tempfile.mkdtemp()) / f"processed.png"
        bg.save(tmp_file)
        return str(tmp_file)
    elif img.mode == 'P':
        img = img.convert('RGBA')
        bg = Image.new('RGB', img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[3])
        tmp_file = Path(tempfile.mkdtemp()) / f"processed.png"
        bg.save(tmp_file)
        return str(tmp_file)
    else:
        return None  # No processing needed


def extract_and_create_ppt(pptx_path, output_pptx):
    """Extract arch images and create a new PPT with white background."""
    pptx_path = Path(pptx_path)
    output_pptx = Path(output_pptx)

    with zipfile.ZipFile(pptx_path, 'r') as z:
        shapes = parse_slide1_images(pptx_path)

    bg_shapes = [s for s in shapes if s['is_background']]
    arch_shapes = [s for s in shapes if not s['is_background']]

    print(f"\n{'='*60}")
    print(f"File: {pptx_path.name}")
    print(f"Slide 1 total image shapes: {len(shapes)}")
    print(f"  Background: {len(bg_shapes)}")
    print(f"  Architecture: {len(arch_shapes)}")
    print(f"{'='*60}")

    for s in shapes:
        tag = "[BG]" if s['is_background'] else "[ARCH]"
        print(f"  {tag} {s['internal_name']} "
              f"pos=({s['x']},{s['y']}) EMU "
              f"size=({s['cx']},{s['cy']}) "
              f"coverage={s['coverage']:.1%}")

    if not arch_shapes:
        print("\nWARNING: No architecture images found!")
        return

    # Create new PPT
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add white background rectangle
    bg_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0),
        prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background()

    print(f"\n[OK] Created new PPT with white background")

    # Add arch images
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for i, s in enumerate(arch_shapes):
            internal_name = s['internal_name']
            media_path = f"ppt/media/{internal_name}"
            try:
                img_data = z.read(media_path)
            except KeyError:
                print(f"  [ERROR] Image not found: {media_path}")
                continue

            # Calculate center position
            img_width = Emu(s['cx'])
            img_height = Emu(s['cy'])
            left = Emu((12192000 - s['cx']) / 2)
            top = Emu((6858000 - s['cy']) / 2)

            # Process image
            actual_w, actual_h = None, None
            tmp_path = process_image_for_white_bg(img_data)

            if tmp_path:
                try:
                    with Image.open(tmp_path) as img:
                        actual_w, actual_h = img.size
                except Exception:
                    pass

            slide.shapes.add_picture(tmp_path or BytesIO(img_data), left, top, img_width, img_height)

            print(f"  [OK] {internal_name}: "
                  f"actual={actual_w}x{actual_h} or {s['cx']}x{s['cy']} EMU, "
                  f"pos=({left/914400:.1f},{top/914400:.1f})in")

            # Cleanup
            if tmp_path and Path(tmp_path).exists():
                try:
                    Path(tmp_path).unlink()
                except Exception:
                    pass

    # Save
    prs.save(str(output_pptx))
    print(f"\n{'='*60}")
    print(f"DONE! Output: {output_pptx}")
    print(f"{'='*60}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) >= 3 else "arch_output.pptx"

    if not Path(pptx_path).exists():
        print(f"[ERROR] File not found: {pptx_path}")
        sys.exit(1)

    extract_and_create_ppt(pptx_path, output_pptx)
