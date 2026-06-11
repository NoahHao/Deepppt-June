#!/usr/bin/env python3
"""
extract_cluster.py — 智能提取 PPTX 第一页的"架构图"元素组
============================================================
核心思路:
  1. 解析 slide1.xml，提取所有 shape 的视觉边界框 (x, y, x+cx, y+cy)
  2. 将所有边界框做 DBSCAN 聚类，找出密集的元素组
  3. 取最大的簇 → 这就是用户说的"架构图"
  4. 将所有元素绘制到一张图片上（白色底色）
  5. 将图片居中放置到新 PPTX 一页中输出

用法:
  python extract_cluster.py <pptx_path> [output_dir]
"""

import os
import re
import sys
import zipfile
from pathlib import Path
from io import BytesIO
from collections import defaultdict

try:
    from PIL import Image
    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False

# ── EMU → 像素 ──
EMU_PER_PX = 914400 / 96

# 幻灯片标准尺寸 (16:9)
SLIDE_EMU_W = 12192000
SLIDE_EMU_H = 6858000

# DBSCAN 参数: 两个元素中心点 Manhattan 距离 < 这个值视为同一簇
CLUSTER_EPS_EMU = 800000  # ~83px at 96dpi


def emu_to_px(emu_val):
    return int(emu_val / EMU_PER_PX)


def parse_slide1_elements(pptx_path):
    """
    解析 slide1.xml，提取所有 shape 的视觉边界框。
    排除全页背景图。
    """
    with zipfile.ZipFile(pptx_path, 'r') as z:
        slide_xml = z.read('ppt/slides/slide1.xml').decode('utf-8', errors='replace')
        try:
            rels_xml = z.read('ppt/slides/_rels/slide1.xml.rels').decode('utf-8', errors='replace')
        except Exception:
            rels_xml = ""

    elements = []
    tag_starts = list(re.finditer(r'<(p:sp|p:pic|p:graphicFrame|p:grpSp)\b', slide_xml))

    for start_m in tag_starts:
        tag_name = start_m.group(1)
        start_pos = start_m.start()

        # 追踪深度找闭合标签
        depth = 1
        search_pos = start_m.end()
        end_pos = -1
        open_tag = f'<{tag_name}'
        close_tag = f'</{tag_name}>'

        while depth > 0 and search_pos < len(slide_xml):
            next_open = slide_xml.find(open_tag, search_pos)
            next_close = slide_xml.find(close_tag, search_pos)
            if next_close == -1:
                break
            if next_open != -1 and next_open < next_close:
                depth += 1
                search_pos = next_open + len(open_tag)
            else:
                depth -= 1
                search_pos = next_close + len(close_tag)
                if depth == 0:
                    end_pos = next_close + len(close_tag)

        if end_pos == -1:
            continue

        block = slide_xml[start_pos:end_pos]

        # 提取名称
        name_m = re.search(r'<p:cNvPr[^>]*name="([^"]*)"', block)
        name = name_m.group(1) if name_m else "?"

        # 提取 spPr 中的 xfrm
        spPr_match = re.search(r'<p:spPr\b[^>]*>.*?</p:spPr>', block, re.DOTALL)
        if not spPr_match:
            spPr_match = re.search(r'<p:grpSpPr\b[^>]*>.*?</p:grpSpPr>', block, re.DOTALL)
        if not spPr_match:
            continue

        xfrm_match = re.search(r'<a:xfrm\b.*?</a:xfrm>', spPr_match.group(), re.DOTALL)
        if not xfrm_match:
            continue

        xfrm_text = xfrm_match.group()
        off_m = re.search(r'<a:off[^>]*x="(\d+)"[^>]*y="(\d+)"', xfrm_text)
        ext_m = re.search(r'<a:ext[^>]*cx="(\d+)"[^>]*cy="(\d+)"', xfrm_text)
        if not off_m or not ext_m:
            continue

        x = int(off_m.group(1))
        y = int(off_m.group(2))
        cx = int(ext_m.group(1))
        cy = int(ext_m.group(2))

        # 跳过空白元素
        if cx <= 0 or cy <= 0:
            continue

        # 跳过全页背景
        if cx >= 12000000 and cy >= 6800000:
            continue

        # 跳过极小装饰点 (< 50px)
        if emu_to_px(cx) < 50 and emu_to_px(cy) < 50:
            continue

        # 判断是否有图片嵌入
        is_image = '<a:blip' in block

        elements.append({
            'name': name,
            'tag': tag_name,
            'x': x, 'y': y,
            'cx': cx, 'cy': cy,
            'x2': x + cx,
            'y2': y + cy,
            'is_image': is_image,
            'area': cx * cy,
            'block': block,
        })

    return elements


def cluster_elements(elements):
    """
    基于 DBSCAN 算法聚类元素。
    两个元素中心点 Manhattan 距离 < CLUSTER_EPS_EMU，视为同一簇。
    返回所有簇，按总面积降序排序。
    """
    centers = []
    for i, e in enumerate(elements):
        cx = (e['x'] + e['x2']) / 2
        cy = (e['y'] + e['y2']) / 2
        centers.append((cx, cy, i))

    visited = set()
    clusters = []

    for i, (cx1, cy1, idx1) in enumerate(centers):
        if i in visited:
            continue

        cluster_indices = {idx1}
        queue = [idx1]
        visited.add(i)

        while queue:
            ci = queue.pop(0)
            cx2, cy2, _ = centers[ci]

            for j, (cx3, cy3, idx2) in enumerate(centers):
                if j in visited:
                    continue
                dist = abs(cx2 - cx3) + abs(cy2 - cy3)
                if dist < CLUSTER_EPS_EMU:
                    cluster_indices.add(idx2)
                    visited.add(j)
                    queue.append(idx2)

        clusters.append(cluster_indices)

    # 按簇的总面积降序排序
    clusters.sort(key=lambda ci: sum(elements[i]['area'] for i in ci), reverse=True)
    return clusters


def extract_cluster_to_pptx(pptx_path, cluster_indices, output_dir):
    """
    将元素组提取为一张图片，居中放置到新 PPTX 一页中输出。
    """
    if not HAS_PILLOW:
        print("[ERROR] Pillow not installed. Install: pip install Pillow")
        return None

    # 重新从文件解析获取完整元素列表
    all_elements = parse_slide1_elements(pptx_path)
    cluster_elements_list = [all_elements[i] for i in cluster_indices]

    # 计算包围盒
    min_x = min(e['x'] for e in cluster_elements_list)
    min_y = min(e['y'] for e in cluster_elements_list)
    max_x = max(e['x2'] for e in cluster_elements_list)
    max_y = max(e['y2'] for e in cluster_elements_list)

    # 内边距
    padding_x = int(SLIDE_EMU_W * 0.03)
    padding_y = int(SLIDE_EMU_H * 0.03)

    bbox_w = max_x - min_x
    bbox_h = max_y - min_y
    min_size = 300 * EMU_PER_PX
    if bbox_w < min_size:
        min_x -= (min_size - bbox_w) // 2
    if bbox_h < min_size:
        min_y -= (min_size - bbox_h) // 2

    width = max_x - min_x + padding_x * 2
    height = max_y - min_y + padding_y * 2

    if width <= 0 or height <= 0:
        print("[ERROR] Invalid bounding box")
        return None

    # 创建白色背景
    img = Image.new('RGB', (width, height), (255, 255, 255))

    # 解析 rels
    with zipfile.ZipFile(pptx_path, 'r') as z:
        rels_xml = z.read('ppt/slides/_rels/slide1.xml.rels').decode('utf-8', errors='replace')
        # 建立 embed_id -> Target 映射
        rel_map = {}
        for m in re.finditer(r'Id="([^"]+)"[^>]*Target="([^"]+)"', rels_xml):
            rel_map[m.group(1)] = m.group(2)

        for e in cluster_elements_list:
            rx = e['x'] - min_x + padding_x
            ry = e['y'] - min_y + padding_y
            rw = e['cx']
            rh = e['cy']

            if e['is_image']:
                try:
                    blip_m = re.search(r'<a:blip[^>]*r:embed="([^"]+)"', e['block'])
                    if blip_m:
                        embed_id = blip_m.group(1)
                        target = rel_map.get(embed_id, '')
                        media_name = Path(target).name
                        img_data = z.read(f"ppt/media/{media_name}")
                        pic = Image.open(BytesIO(img_data))

                        target_w = emu_to_px(rw)
                        target_h = emu_to_px(rh)
                        if pic.size[0] != target_w or pic.size[1] != target_h:
                            pic = pic.resize((target_w, target_h), Image.LANCZOS)

                        if pic.mode == 'RGBA':
                            bg = Image.new('RGB', pic.size, (255, 255, 255))
                            bg.paste(pic, mask=pic.split()[3])
                            pic = bg
                        elif pic.mode == 'P':
                            pic = pic.convert('RGBA')
                            bg = Image.new('RGB', pic.size, (255, 255, 255))
                            bg.paste(pic, mask=pic.split()[3])
                            pic = bg

                        img.paste(pic, (rx, ry))
                except Exception as err:
                    print(f"  [WARN] Image failed: {e['name']} - {err}")

            elif e['tag'] == 'p:sp':
                # 文本框/形状: 尝试提取填充色
                fill_m = re.search(r'<a:solidFill><a:srgbClr val="([0-9a-fA-F]+)"', e['block'])
                rect_w = emu_to_px(rw)
                rect_h = emu_to_px(rh)

                if fill_m:
                    color_hex = fill_m.group(1)
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    color = (r, g, b)
                    for y_off in range(rect_h):
                        for x_off in range(rect_w):
                            img.putpixel((rx + x_off, ry + y_off), color)

                    # 边框
                    stroke_m = re.search(r'<a:ln[^>]*w="(\d+)">', e['block'])
                    if stroke_m:
                        stroke_m2 = re.search(r'<a:solidFill><a:srgbClr val="([0-9a-fA-F]+)"', e['block'])
                        if stroke_m2:
                            sc = stroke_m2.group(1)
                            sr, sg, sb = int(sc[0:2], 16), int(sc[2:4], 16), int(sc[4:6], 16)
                            for x_off in range(rect_w):
                                img.putpixel((rx + x_off, ry), (sr, sg, sb))
                                img.putpixel((rx + x_off, ry + rect_h - 1), (sr, sg, sb))
                            for y_off in range(rect_h):
                                img.putpixel((rx, ry + y_off), (sr, sg, sb))
                                img.putpixel((rx + rect_w - 1, ry + y_off), (sr, sg, sb))

    # 保存图片
    tmp_img = output_dir / "extracted_cluster.png"
    img.save(tmp_img)

    # 创建 PPTX
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("[WARN] python-pptx not available, only saved PNG")
        return str(tmp_img)

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # 空白布局
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 白色背景
    background = slide.background
    fill = background.fill
    fill.solid()

    # 居中放置图片
    left = (12192000 - img.width * EMU_PER_PX) // 2
    top = (6858000 - img.height * EMU_PER_PX) // 2

    slide.shapes.add_picture(
        str(tmp_img),
        left=Emu(left),
        top=Emu(top),
        width=img.width * EMU_PER_PX,
        height=img.height * EMU_PER_PX
    )

    out_path = output_dir / "extracted_architecture.pptx"
    prs.save(str(out_path))
    return str(out_path)


# ═══════════ CLI ═══════════
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = Path(sys.argv[2]) if len(sys.argv) >= 3 else Path(__file__).parent / "cluster_output"

    if not os.path.exists(pptx_path):
        print(f"[ERROR] File not found: {pptx_path}")
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"\n{'='*60}")
    print(f"File: {Path(pptx_path).name}")
    print(f"{'='*60}")

    # 1. 解析
    elements = parse_slide1_elements(pptx_path)
    print(f"\nFound {len(elements)} non-background elements on slide 1")

    for e in elements:
        tag = "[IMG]" if e['is_image'] else "[TXT]"
        print(f"  {tag} {e['name']:20s} "
              f"pos=({emu_to_px(e['x'])},{emu_to_px(e['y'])}) "
              f"size={emu_to_px(e['cx'])}x{emu_to_px(e['cy'])}px "
              f"center=({emu_to_px((e['x']+e['x2'])/2)},{emu_to_px((e['y']+e['y2'])/2)})")

    # 2. 聚类
    clusters = cluster_elements(elements)
    print(f"\nClustered into {len(clusters)} groups:")

    for i, ci in enumerate(clusters):
        elements_in_cluster = [elements[idx] for idx in ci]
        total_area = sum(e['area'] for e in elements_in_cluster)
        print(f"\n  Cluster {i+1}: {len(ci)} elements, area={emu_to_px(int(total_area**0.5))}x{emu_to_px(int(total_area**0.5))}px")
        for idx in ci:
            e = elements[idx]
            tag = "[IMG]" if e['is_image'] else "[TXT]"
            print(f"    {tag} {e['name']:20s} center=({emu_to_px((e['x']+e['x2'])/2)},{emu_to_px((e['y']+e['y2'])/2)})px")

    # 3. 取最大簇
    best_cluster_indices = list(clusters[0]) if clusters else []
    best_cluster = [elements[idx] for idx in best_cluster_indices]

    print(f"\n{'='*60}")
    print(f"Selected largest cluster: {len(best_cluster)} elements")
    print(f"{'='*60}")

    # 4. 输出
    print(f"\nExtracting...")
    result = extract_cluster_to_pptx(pptx_path, best_cluster_indices, output_dir)
    if result:
        print(f"\n{'='*60}")
        print(f"DONE! Output: {result}")
        print(f"{'='*60}")
