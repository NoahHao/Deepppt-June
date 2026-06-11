#!/usr/bin/env python3
"""
PPT 生成模块 - 基于索引图片创建 PPTX

功能：
1. 从 image_extract 索引 JSON 读取图片
2. 等比例缩放图片（防止变形）
3. 纯白背景（无底板/水印）
4. 智能图片展示区域划分
"""

import json
import sys
import os
from pathlib import Path
from typing import List, Optional, Tuple
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ───── 常量 ─────
SLIDE_WIDTH = Inches(13.333)  # 宽屏 16:9
SLIDE_HEIGHT = Inches(7.5)

DEFAULT_TITLE_FONT_SIZE = Pt(28)
DEFAULT_SUBTITLE_FONT_SIZE = Pt(16)
DEFAULT_TEXT_FONT_SIZE = Pt(12)

# 白色 RGB
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)

# 图片排列模式
MODE_COLUMN = "column"   # 竖直排列
MODE_HORIZONTAL = "horizontal"  # 横向排列（默认）


def _load_index(index_path: str) -> dict:
    """加载 JSON 索引文件"""
    with open(index_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def _calc_fit_image(
    img_width: int, img_height: int,
    box_width: Emu, box_height: Emu,
) -> Tuple[Emu, Emu, Emu, Emu]:
    """
    等比例缩放图片，使其完全填充到指定矩形区域内，不裁剪，不拉伸变形。
    返回 (x, y, w, h) 单位 EMU，居中对齐在区域内。
    """
    # 计算原始宽高比
    aspect_ratio = img_width / max(img_height, 1)

    # 将 box 宽高转换为 EMU（它们已是 Emu 类型）
    bw = box_width
    bh = box_height

    # 按宽约束
    w_from_width = bw
    h_from_width = int(bw / aspect_ratio)

    # 按高约束
    h_from_height = bh
    w_from_height = int(bh * aspect_ratio)

    if h_from_width <= bh:
        # 以宽度为准，高度不溢出
        final_w = w_from_width
        final_h = h_from_width
    else:
        # 以高度为准，宽度不溢出
        final_w = w_from_height
        final_h = h_from_height

    # 居中偏移
    x = int((bw - final_w) // 2)
    y = int((bh - final_h) // 2)

    return x, y, final_w, final_h


def _add_white_background(slide, width_emu=SLIDE_WIDTH, height_emu=SLIDE_HEIGHT):
    """
    为幻灯片添加纯白矩形背景，覆盖整张幻灯片。
    注意：纯白背景覆盖原模板所有底板/水印。
    """
    # 插入一个纯白矩形，位于最底层（通过先插入实现）
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width_emu, height_emu
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE_RGB
    shape.line.fill.background()  # 无边框
    # 将形状发送到最底层
    # python-pptx 没有直接 send_to_back，但调整 z-order 可用
    # 通过 move 到第一个位置（在 XML 层面）
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)  # 插入到背景之后、内容之前
    return shape


def _arrange_images(
    images: List[dict],
    slide_width: Emu,
    slide_height: Emu,
    margin_x: Emu = Inches(0.5),
    margin_y: Emu = Inches(0.5),
    gap: Emu = Inches(0.2),
    title_height: Emu = Inches(1.0),
) -> List[Tuple[Emu, Emu, Emu, Emu]]:
    """
    智能划分图片展示区域。

    支持模式：
    - 1 张：居中展示，最大尺寸
    - 2 张：左右平分
    - 3 张：上 1 下 2
    - 4 张及以上：网格排列（每行 ≤3）

    返回 [(x, y, w, h), ...] 框坐标（EMU）
    """
    n = len(images)

    usable_w = slide_width - 2 * margin_x
    usable_h = slide_height - margin_y - title_height - margin_y

    if n <= 1:
        # 居中，最大尺寸（留边距）
        x = margin_x
        y = title_height + margin_y
        w = usable_w
        h = usable_h
        return [(x, y, w, h)]

    elif n == 2:
        # 左右平分
        half_w = (usable_w - gap) // 2
        y = title_height + margin_y
        return [
            (margin_x, y, half_w, usable_h),
            (margin_x + half_w + gap, y, half_w, usable_h),
        ]

    elif n == 3:
        # 上 1 下 2
        top_h = (usable_h - gap) // 2
        bottom_h = top_h
        half_w = (usable_w - gap) // 2
        # 上
        top_w = usable_w
        top_x = margin_x
        top_y = title_height + margin_y
        # 下左
        bottom_left_x = margin_x
        bottom_left_y = top_y + top_h + gap
        bottom_right_x = margin_x + half_w + gap
        bottom_right_y = bottom_left_y
        return [
            (top_x, top_y, top_w, top_h),
            (bottom_left_x, bottom_left_y, half_w, bottom_h),
            (bottom_right_x, bottom_right_y, half_w, bottom_h),
        ]

    else:  # n >= 4
        # 网格排列，每行最多 3 列
        cols = min(3, n)
        rows = (n + cols - 1) // cols
        cell_w = (usable_w - (cols - 1) * gap) // cols
        cell_h = (usable_h - (rows - 1) * gap) // rows

        boxes = []
        for idx in range(n):
            row = idx // cols
            col = idx % cols
            x = margin_x + col * (cell_w + gap)
            y = title_height + margin_y + row * (cell_h + gap)
            boxes.append((x, y, cell_w, cell_h))
        return boxes


def _add_title_text(slide, text: str, slide_width_emu: Emu):
    """在幻灯片顶部添加标题文字"""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), slide_width_emu - Inches(1.0), Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = DEFAULT_TITLE_FONT_SIZE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT


def _guess_title(image_entries: list) -> str:
    """
    根据图片集合智能生成幻灯片标题。
    优先使用公共来源文件名，否则用当前目录名。
    """
    dir_names = set()
    base_names = set()
    for entry in image_entries:
        src_path = entry.get("source_file", "")
        if src_path:
            parts = Path(src_path).parts
            if len(parts) >= 2:
                dir_names.add(parts[-2])
            base_names.add(Path(src_path).stem)
    if len(dir_names) == 1:
        return list(dir_names)[0]
    if base_names:
        return next(iter(base_names))
    return "图片展示"


def generate_ppt(
    index_path: str,
    output_path: str,
    kb_root: Optional[str] = None,
    max_ppt_images: int = 6,
) -> str:
    """
    从图片索引生成 PPT。

    Args:
        index_path: JSON 索引文件路径
        output_path: 输出的 PPTX 文件路径
        kb_root: 可选，知识库根目录（用于筛选图片）
        max_ppt_images: 每张幻灯片最大图片数（默认 6）

    Returns:
        str: 生成的 PPTX 文件路径
    """
    data = _load_index(index_path)
    images = data.get("images", {})
    if not images:
        raise ValueError("索引中无图片数据")

    # 按来源文件分组
    file_groups = defaultdict(list)
    for key, entry in images.items():
        src = entry.get("source_file", "unknown")
        file_groups[src].append(entry)

    # 创建 Presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # 遍历每个文件组，每组生成一页或多页
    for file_key, img_entries in file_groups.items():
        # 按 slide_index 排序
        img_entries.sort(key=lambda e: e.get("slide_index", 0))

        # 每页 max_ppt_images 张
        for chunk_start in range(0, len(img_entries), max_ppt_images):
            chunk = img_entries[chunk_start:chunk_start + max_ppt_images]

            slide = prs.slides.add_slide(blank_layout)

            # 添加纯白背景
            _add_white_background(slide)

            # 添加标题
            slide_title = _guess_title(chunk)
            _add_title_text(slide, slide_title, SLIDE_WIDTH)

            # 计算图片框位置
            boxes = _arrange_images(chunk, SLIDE_WIDTH, SLIDE_HEIGHT)

            for idx, entry in enumerate(chunk):
                if idx >= len(boxes):
                    break

                img_path = entry.get("archive_path", "")
                if not img_path or not os.path.isfile(img_path):
                    continue

                img_w = entry.get("width", 0) or entry.get("size_width", 0)
                img_h = entry.get("height", 0) or entry.get("size_height", 0)

                box_x, box_y, box_w, box_h = boxes[idx]

                # 等比例缩放并居中
                if img_w > 0 and img_h > 0:
                    ix, iy, iw, ih = _calc_fit_image(img_w, img_h, box_w, box_h)
                    final_x = box_x + ix
                    final_y = box_y + iy
                    final_w = iw
                    final_h = ih
                else:
                    final_x, final_y, final_w, final_h = box_x, box_y, box_w, box_h

                slide.shapes.add_picture(img_path, final_x, final_y, final_w, final_h)

    # 保存
    output_path_obj = Path(output_path)
    output_path_obj.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path_obj))
    return str(output_path_obj)


def main():
    """CLI 入口"""
    if len(sys.argv) < 3:
        print("用法:")
        print("  python ppt_generator.py <index_json_path> <output_pptx_path>")
        print("选项:")
        print("  --max-images N    每张幻灯片最大图片数 (默认: 6)")
        sys.exit(1)

    index_path = sys.argv[1]
    output_path = sys.argv[2]

    max_images = 6
    if "--max-images" in sys.argv:
        idx = sys.argv.index("--max-images")
        if idx + 1 < len(sys.argv):
            max_images = int(sys.argv[idx + 1])

    generate_ppt(index_path, output_path, max_ppt_images=max_images)
    print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    main()
