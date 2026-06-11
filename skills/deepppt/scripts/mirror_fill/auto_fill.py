#!/usr/bin/env python3
"""
mirror_fill auto_fill — 一键检索架构图 + 自动填充
====================================================

固化搜索→分组→映射→填充全流程。用户只需提供两个输入：
  1. query       — 搜索查询，定位源 PPT 页
  2. new_blocks  — 要替换的新内容（每个 block 一行或多行文字）

系统自动：
  - 调用 LLM_Search 定位最优匹配的 PPT 页
  - 分析目标区域文字结构的段落分组
  - 将 new_blocks 按段落逐一映射到旧文字
  - 在 XML 层完成替换，保持格式不变

用法：
    from mirror_fill.auto_fill import mirror_auto
    
    mirror_auto(
        "DCS数据中心解决方案架构图",
        [
            "企业级高可用\\n"
            "vSphere HA 自动故障切换，vMotion 零停机迁移\\n"
            "历经 25 年全球企业级验证",
            
            "智能统一运维\\n"
            "vCenter 集中管控，vRealize AIOps 智能运维\\n"
            "多云环境统一视图，自动化策略管理",
            
            "开放生态兼容\\n"
            "500+ ISV 认证，Tanzu 容器平台\\n"
            "混合云无缝迁移扩展，vSAN 超融合架构",
        ]
    )
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

THIS_DIR = Path(__file__).resolve().parent
SCRIPTS_DIR = THIS_DIR.parent
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_DIR))

from pptx.util import Emu

# 兼容 CLI 直接运行 (python auto_fill.py) 和包导入 (from mirror_fill.auto_fill import ...)
try:
    from .filler import MirrorFiller, _find_text_shapes_in_region
    from .regions import get_layout, SlideRegion, PRESET_LAYOUTS
except ImportError:
    from mirror_fill.filler import MirrorFiller, _find_text_shapes_in_region
    from mirror_fill.regions import get_layout, SlideRegion, PRESET_LAYOUTS

# ═══════════════════════════════════════════════════════
# 自动分组引擎
# ═══════════════════════════════════════════════════════

# top-to-top 间距候选阈值（EMU）
# slide#1: 组内 ~550k, 组间 ~1.3M → 800k 合适
# slide#3: 组内 ~350k, 组间 ~1.4M → 800k 合适
GAP_CANDIDATES = [800_000, 600_000, 1_000_000, 400_000, 1_200_000]


def _auto_select_gap_threshold(
    slide: "pptx.slide.Slide",
    shapes: List[Tuple[int, str]],
    preferred_n_blocks: int = 0,
) -> int:
    """自动选择最优分组间隔阈值。

    尝试多个候选值，选择最自然的 block 数（2-5 为佳）。
    preferred_n_blocks 仅作为加分项，不会导致极端分组（1或10+）。
    """
    best_gap = GAP_CANDIDATES[0]
    best_score = -1

    for gap in GAP_CANDIDATES:
        groups = _group_region_shapes(slide, shapes, gap_threshold=gap)
        n = len(groups)

        if n < 1:
            continue

        # 基础分: 2-5 blocks 为最佳区间
        if 2 <= n <= 5:
            base_score = 20 - abs(n - 3) * 3  # 3 blocks = 20分, 2/4 = 17分, 5 = 14分
        elif 6 <= n <= 8:
            base_score = 10 - n  # 勉强可用
        else:
            base_score = 0  # 太极端

        # 匹配加分: 与用户预期一致
        match_bonus = 5 if (0 < preferred_n_blocks == n) else 0

        score = base_score + match_bonus
        if score > best_score:
            best_score = score
            best_gap = gap

    return best_gap


def _group_region_shapes(
    slide: "pptx.slide.Slide",
    shapes: List[Tuple[int, str]],
    gap_threshold: int = 400_000,
) -> List[List[Tuple[int, str]]]:
    """按 Y 坐标间距将区域内的 shape 分组为逻辑 blocks。

    每组对应一个逻辑单元（如一个优势项 = 标题 + 描述）。

    Args:
        slide: python-pptx Slide 对象
        shapes: [(shape_index, text), ...]，来自 _find_text_shapes_in_region
        gap_threshold: 分组间隔阈值 (EMU)，大于此值认为新 block

    Returns:
        [[(idx, text), ...], ...]  — 每组含该 block 中的所有 shape
    """
    if not shapes:
        return []

    # 按 Y 坐标排序
    sorted_shapes = sorted(
        shapes,
        key=lambda s: slide.shapes[s[0]].top,
    )

    groups: List[List[Tuple[int, str]]] = []
    current_group: List[Tuple[int, str]] = [sorted_shapes[0]]

    for i in range(1, len(sorted_shapes)):
        prev_idx, _ = sorted_shapes[i - 1]
        curr_idx, _ = sorted_shapes[i]

        prev_shape = slide.shapes[prev_idx]
        curr_shape = slide.shapes[curr_idx]

        # 使用 top-to-top 间距 (避免 shape 重叠导致 bottom-based 失效)
        gap = curr_shape.top - prev_shape.top

        if gap > gap_threshold:
            groups.append(current_group)
            current_group = []

        current_group.append(sorted_shapes[i])

    if current_group:
        groups.append(current_group)

    return groups


def _extract_paragraphs_from_group(
    slide: "pptx.slide.Slide",
    group: List[Tuple[int, str]],
) -> List[str]:
    """从一组 shape 中按顺序提取所有段落文本。

    遍历 group 中的每个 shape，按段落顺序收集文本。
    跳过空段落。
    """
    paras: List[str] = []
    for idx, _ in group:
        shape = slide.shapes[idx]
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                paras.append(text)
    return paras


def _build_text_map_from_blocks(
    slide: "pptx.slide.Slide",
    groups: List[List[Tuple[int, str]]],
    new_blocks: List[str],
    *,
    auto_balance: bool = True,
    verbose: bool = True,
) -> Dict[str, str]:
    """将自动分组的旧段落与用户提供的新内容块一一映射。

    支持 auto_balance 模式：
      - 用户块数 > 检测块数 → 自动聚类合并 (\n\n 为分隔)
      - 用户块数 < 检测块数 → 自动拆分或用最后块填充
      - 段落数不匹配 → 打印提示，超额段落附加到最后一行

    Args:
        slide: Slide 对象
        groups: _group_region_shapes 的输出
        new_blocks: 用户提供的新内容列表
        auto_balance: 是否自动平衡块数

    Returns:
        {old_text: new_text} 映射字典
    """
    n_groups = len(groups)
    n_blocks = len(new_blocks)

    # ── 自动平衡块数 ──
    if auto_balance and n_groups != n_blocks:
        if n_blocks > n_groups:
            # 用户块数过多 → 聚类合并
            if verbose:
                print(f"  ⚠️ 用户提供 {n_blocks} 块，需要 {n_groups} 块，自动聚类合并")
            # 按比例分组：每 ceil(n_blocks / n_groups) 个用户块合并为1个
            chunk_size = (n_blocks + n_groups - 1) // n_groups
            merged = []
            for gi in range(n_groups):
                chunk = new_blocks[gi * chunk_size:(gi + 1) * chunk_size]
                # 用 \n\n 连接多个用户块
                merged.append("\n\n".join(chunk))
            new_blocks = merged
        else:
            # 用户块数不足 → 用最后块填充
            if verbose:
                print(f"  ⚠️ 用户仅提供 {n_blocks} 块，需要 {n_groups} 块，用最后块填充")
            while len(new_blocks) < n_groups:
                new_blocks.append(new_blocks[-1])

    # ── 生成映射 ──
    text_map: Dict[str, str] = {}
    all_hints = []

    for gi, group in enumerate(groups):
        old_paras = _extract_paragraphs_from_group(slide, group)

        new_block_text = new_blocks[gi].replace('\\n', '\n')
        new_paras = [p.strip() for p in new_block_text.split('\n') if p.strip()]

        n_old = len(old_paras)
        n_new = len(new_paras)

        if n_new < n_old and verbose:
            all_hints.append(f"    Block {gi+1}: 需要 {n_old} 段，你提供了 {n_new} 段")
        elif n_new > n_old:
            # 超额段落合并到最后一段
            extra = "；".join(new_paras[n_old - 1:])
            new_paras = new_paras[:n_old - 1] + [extra]
            if verbose:
                all_hints.append(f"    Block {gi+1}: 需要 {n_old} 段，你提供了 {n_new} 段（已合并超额内容）")

        n_map = min(n_old, len(new_paras))
        for pi in range(n_map):
            old_text = old_paras[pi]
            new_text = new_paras[pi]
            if old_text != new_text:
                text_map[old_text] = new_text

    # ── 打印段落提示 ──
    if all_hints and verbose:
        print("  💡 段落数量提示:")
        for h in all_hints:
            print(h)
        print(f"     每个 block 需要填写的段落结构一览：")
        for gi, group in enumerate(groups):
            paras = _extract_paragraphs_from_group(slide, group)
            print(f"       Block {gi+1}: {len(paras)} 段 — [标题] + {len(paras)-1} 行描述")

    return text_map


# ═══════════════════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════════════════

def mirror_auto(
    query: str,
    new_blocks: List[str],
    *,
    region: str = "right",
    layout: str = "left_right",
    output_dir: str | Path | None = None,
    output_name: str | None = None,
    dry_run: bool = False,
    verbose: bool = True,
) -> Path | None:
    """一键检索 PPT 页面并自动填充文字。

    完整流程:
      1. LLM_Search 语义搜索 → 定位最佳匹配的 PPT 页
      2. 读取目标区域 (默认 right) 的文字结构
      3. 自动按 Y 坐标间距分组为逻辑 blocks
      4. 将 new_blocks 与检测到的 blocks 逐段落映射
      5. XML 层执行替换（保持格式不变）
      6. 输出到指定目录

    Args:
        query: 搜索查询，如 "DCS数据中心解决方案架构图"
        new_blocks: 新内容列表，每个元素为一个逻辑块的完整文字。
                    块内各段落用 \\n 分隔。
                    例如: ["标题1\\n描述行1\\n描述行2", "标题2\\n描述行1\\n描述行2", ...]
        region: 目标区域名，默认 "right"
        layout: 布局预设名，默认 "left_right"
        output_dir: 输出目录，默认 deepppt/project/
        output_name: 输出文件名，默认自动生成
        dry_run: True 时仅分析不保存
        verbose: True 时打印详细信息

    Returns:
        输出 PPTX 路径（dry_run 时返回 None）
    """
    # ── Step 1: 搜索 ──
    if verbose:
        print("=" * 60)
        print(f"  🔍 搜索: {query}")
        print("=" * 60)

    filler = MirrorFiller.from_search(query, layout="left_right")  # 初始布局，后续自动探测

    if verbose:
        print(f"  定位: {filler.src_pptx.name}  slide #{filler.slide_num}")
        print()

    # ── Step 2 & 3: 读区域 + 自动分组（智能回退多布局）──
    # 评分策略：优先匹配 block 数与用户输入一致，其次选 block 内段落数均匀的
    candidate_layouts = [
        (layout, region),                      # 用户指定的布局+区域
        ("left_center_right", "center"),       # 左中右 center（常见架构图正文区）
        ("left_center_right", "right"),        # 左中右 right
        (layout, "left"),                      # 同布局 left
        (layout, "center") if hasattr(get_layout(layout), 'center') else None,
        ("left_center_right", "left"),         # 左中右 left
        ("top_bottom", "bottom"),              # 上下布局
    ]
    candidate_layouts = [c for c in candidate_layouts if c is not None]

    best_used_region = None
    best_used_layout = None
    best_shapes = None
    best_groups = None
    best_layout_regions = None
    best_score = -1  # 评分：block数匹配 + 段落均匀度

    for try_layout, try_region in candidate_layouts:
        try:
            try_regions = get_layout(try_layout)
            if try_region not in try_regions:
                continue
            target = try_regions[try_region]
            shapes = _find_text_shapes_in_region(filler._slide, target)
            if not shapes:
                continue

            # 动态选择最优分组阈值（不与用户输入绑定）
            best_gap = _auto_select_gap_threshold(filler._slide, shapes, 0)
            groups = _group_region_shapes(filler._slide, shapes, gap_threshold=best_gap)
            if not groups or len(groups) < 2:
                continue

            n = len(groups)

            # 计算平均每 block 段落数方差（均匀度）
            para_counts = [len(_extract_paragraphs_from_group(filler._slide, g)) for g in groups]
            avg_paras = sum(para_counts) / n if n else 0
            variance = sum((p - avg_paras) ** 2 for p in para_counts) / n if n else 999

            # 综合评分: 目标 3 blocks（架构图最常见） + 低方差 + 居中偏好
            block_score = 20 - abs(n - 3) * 4           # 3 blocks=20分
            uniformity = max(0, 10 - variance)           # 均匀=10分
            para_score = min(5, avg_paras) if 1.5 <= avg_paras <= 8 else 0
            center_bonus = 3 if try_region == "center" else 0  # center 区域优先

            score = block_score + uniformity + para_score + center_bonus
            if score > best_score:
                best_score = score
                best_used_region = try_region
                best_used_layout = try_layout
                best_shapes = shapes
                best_groups = groups
                best_layout_regions = try_regions
        except (ValueError, KeyError):
            continue

    # 无任何布局匹配
    if not best_shapes:
        # 最终回退: 全幻灯片
        full = SlideRegion(name="full", x_range=None, y_range=None)
        best_shapes = _find_text_shapes_in_region(filler._slide, full)
        if best_shapes:
            best_used_region = "full"
            best_used_layout = layout
            best_groups = _group_region_shapes(filler._slide, best_shapes,
                gap_threshold=_auto_select_gap_threshold(filler._slide, best_shapes, len(new_blocks)))
            best_layout_regions = get_layout(layout)

    if not best_shapes:
        raise ValueError(
            f"所有布局的目标区域均未找到文字 shape。\n"
            f"源页: {filler.src_pptx.name} slide #{filler.slide_num}"
        )

    used_region = best_used_region
    used_layout = best_used_layout
    region_shapes = best_shapes
    groups = best_groups
    regions = best_layout_regions

    if verbose:
        status = "✅" if used_layout == layout and used_region == region else "⚠️ 回退"
        print(f"  📐 布局 [{used_layout}] 区域 [{used_region}] {status}: "
              f"{len(region_shapes)} 个 shape → {len(groups)} 个 block")
        for gi, group in enumerate(groups):
            paras = _extract_paragraphs_from_group(filler._slide, group)
            print(f"    Block {gi+1} ({len(group)} shapes, {len(paras)} 段)")
            for pi, p in enumerate(paras):
                print(f"      [{pi}] {p[:80]}")
        print()

    # ── Step 4: 映射 ──
    text_map = _build_text_map_from_blocks(filler._slide, groups, new_blocks)

    if verbose:
        print(f"  🔗 生成 {len(text_map)} 条映射")
        for old, new in list(text_map.items())[:5]:
            print(f"    {old[:40]:40s} → {new[:40]}")
        if len(text_map) > 5:
            print(f"    ... 还有 {len(text_map) - 5} 条")
        print()

    if dry_run:
        if verbose:
            print("  [dry-run] 预览结束，跳过保存")
        return None

    # ── Step 5: 填充 ──
    # 确保 filler 内部布局名与实际使用的布局一致
    if used_layout != filler._layout_name:
        filler._layout_name = used_layout
        filler._layout_regions = regions

    if used_region in regions:
        filler.fill(used_layout, {used_region: text_map})
    else:
        # 回退区域不在预设布局中，用自定义区域
        custom_layout = {used_region: SlideRegion(name=used_region, x_range=None, y_range=None)}
        filler.fill(custom_layout, {used_region: text_map})

    # ── Step 6: 输出 ──
    if output_dir is None:
        output_dir = SCRIPTS_DIR.parent / "projects"
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if output_name is None:
        # 自动生成文件名：从 query 提取关键词
        safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in query)[:30]
        output_name = f"filled_{safe}.pptx"

    output_path = output_dir / output_name
    filler.save(output_path)

    return output_path


# ═══════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="mirror_auto — 一键检索 PPT 并自动填充",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python auto_fill.py "DCS数据中心架构图" "新内容1" "新内容2" "新内容3"
  python auto_fill.py --query "架构图" --blocks blocks.txt
  python auto_fill.py --query "架构图" --dry-run
        """,
    )
    parser.add_argument("query", help="搜索查询，如 'DCS数据中心架构图'")
    parser.add_argument("blocks", nargs="*", help="新内容块列表")
    parser.add_argument("--region", default="right", help="目标区域 (default: right)")
    parser.add_argument("--layout", default="left_right", help="布局预设 (default: left_right)")
    parser.add_argument("--output-dir", "-d", default=None, help="输出目录")
    parser.add_argument("--output-name", "-o", default=None, help="输出文件名")
    parser.add_argument("--dry-run", "-n", action="store_true", help="仅分析不保存")
    parser.add_argument("--quiet", "-q", action="store_true", help="静默模式")

    args = parser.parse_args()

    mirror_auto(
        query=args.query,
        new_blocks=args.blocks,
        region=args.region,
        layout=args.layout,
        output_dir=args.output_dir,
        output_name=args.output_name,
        dry_run=args.dry_run,
        verbose=not args.quiet,
    )
