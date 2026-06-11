#!/usr/bin/env python3
"""
mirror_fill filler — PPTX 文本替换引擎
========================================

从 PPTX 源文件中复制指定 slide，按区域替换文本，输出新 PPTX。

核心技术: 
  - Slide 复制: python-pptx deepcopy element
  - 文本替换:    zipfile + ElementTree 直接操作 XML (继承自 style_convert)
  - 区域定位:    按 EMU 坐标筛选 shape

为什么用 XML 层替换？
  - GROUP shapes 的文本也能稳定替换
  - 不经过 python-pptx Object API，避免 GROUP 兼容问题
  - 与 style_convert 完全兼容，已验证可靠

用法:
    from mirror_fill.filler import MirrorFiller
    
    filler = MirrorFiller("source.pptx", slide_num=1)
    filler.fill("left_right", {
        "right": {
            "安全可靠": "企业级稳定性",
            "硬件亚健康自动迁移": "vSphere HA 自动故障切换",
        }
    })
    filler.save("output.pptx")
"""

from __future__ import annotations

import copy
import os
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu

from .regions import (
    SlideRegion,
    get_layout,
    custom_layout,
    PRESET_LAYOUTS,
    SLIDE_W,
    SLIDE_H,
)

# ═══════════════════════════════════════════════════════
# OOXML Namespace constants (from style_convert)
# ═══════════════════════════════════════════════════════

NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'


# ═══════════════════════════════════════════════════════
# Shape 定位工具
# ═══════════════════════════════════════════════════════

def _find_text_shapes_in_region(
    slide: "pptx.slide.Slide",
    region: SlideRegion,
) -> List[Tuple[int, str]]:
    """在 PPTX slide 中找到指定区域内的文本 shape。

    返回: [(shape_index, shape_text), ...]
    
    使用 python-pptx 的 shape.left/top 进行区域判定，
    因为 Python API 层面更容易获取坐标。
    """
    matches = []
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if region.contains(shape.left, shape.top):
            matches.append((idx, text))
    return matches


def _find_text_in_slide(
    slide: "pptx.slide.Slide",
    keyword: str,
) -> List[Tuple[int, str]]:
    """在整个 slide 中搜索包含关键词的 shape。

    返回: [(shape_index, full_text), ...]
    """
    matches = []
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if keyword in text:
            matches.append((idx, text))
    return matches


# ═══════════════════════════════════════════════════════
# XML 层文本替换引擎 (继承自 style_convert)
# ═══════════════════════════════════════════════════════

def _normalize(text: str) -> str:
    """归一化文本用于模糊匹配。

    处理全角/半角标点、多余空格等差异。
    """
    # 全角转半角
    result = []
    for c in text:
        code = ord(c)
        if 0xFF01 <= code <= 0xFF5E:
            result.append(chr(code - 0xFEE0))
        elif code == 0x3000:  # 全角空格
            result.append(' ')
        else:
            result.append(c)
    text = ''.join(result)
    # 合并多余空格
    text = ' '.join(text.split())
    return text


def _apply_text_replacements_in_pptx(
    pptx_path: Path,
    slide_xml_name: str,
    text_map: Dict[str, str],
) -> int:
    """在 PPTX 的指定 slide XML 中原地替换文本。

    直接操作 ZIP 内的 slideN.xml，用 ElementTree 改 <a:t> 文本。
    完全不经过 python-pptx Object API，GROUP shapes 的文本也能改。

    Args:
        pptx_path: PPTX 文件路径（原地修改）
        slide_xml_name: 如 "ppt/slides/slide1.xml"
        text_map: {old_text: new_text} 映射

    Returns:
        成功替换的文本数
    """
    applied = 0

    # 读取 → 修改 → 写回（需先读出所有内容再重建）
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        items = {}
        for name in zin.namelist():
            if name == slide_xml_name:
                # 修改此文件
                xml_bytes = zin.read(name)
                root = ET.fromstring(xml_bytes)

                # 在段落级别匹配（处理富文本分段问题）
                for para in root.findall(f'.//{NS_A}p'):
                    # 收集该段落内所有 t 元素的完整文本
                    t_elems = para.findall(f'.//{NS_A}t')
                    full_text = ''.join((t.text or '') for t in t_elems).strip()
                    if not full_text:
                        continue

                    matched_new = None
                    # 策略1: 精确匹配
                    if full_text in text_map:
                        matched_new = text_map[full_text]
                    else:
                        # 策略2: 归一化匹配
                        norm = _normalize(full_text)
                        for old, new in text_map.items():
                            if _normalize(old) == norm:
                                matched_new = new
                                break
                        # 策略3: 子串匹配
                        if not matched_new:
                            for old, new in text_map.items():
                                if len(old) > 4 and old in full_text:
                                    matched_new = full_text.replace(old, new)
                                    break

                    if matched_new:
                        # 替换：清空所有 t 元素，把新文本放到第一个
                        for j, t in enumerate(t_elems):
                            if j == 0:
                                t.text = matched_new
                            else:
                                t.text = ''
                        applied += 1

                new_xml = ET.tostring(root, encoding='UTF-8', xml_declaration=True)
                items[name] = new_xml
            else:
                items[name] = zin.read(name)

    # 重建 ZIP
    pptx_path.unlink()
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in items.items():
            zout.writestr(name, data)

    return applied


def _apply_text_replacements_in_memory(
    slide_bytes: bytes,
    slide_xml_name: str,
    text_map: Dict[str, str],
) -> Tuple[bytes, int]:
    """在内存中的 slide XML bytes 上替换文本。

    用于在保存前批量处理多个 slide XML。

    Returns:
        (modified_bytes, count_of_replacements)
    """
    root = ET.fromstring(slide_bytes)
    applied = 0

    for t_elem in root.findall(f'.//{NS_A}t'):
        text = (t_elem.text or '').strip()
        if not text:
            continue

        if text in text_map:
            t_elem.text = text_map[text]
            applied += 1
            continue

        for old, new in text_map.items():
            if old in text and old != text:
                t_elem.text = text.replace(old, new)
                applied += 1
                break

    return ET.tostring(root, encoding='UTF-8', xml_declaration=True), applied


# ═══════════════════════════════════════════════════════
# MirrorFiller — 主引擎
# ═══════════════════════════════════════════════════════

class MirrorFiller:
    """PPTX 镜像填充引擎。

    典型工作流:
        filler = MirrorFiller.from_pptx("source.pptx", slide_num=1)
        filler.fill("left_right", {
            "right": {"旧文字1": "新文字1", "旧文字2": "新文字2"}
        })
        filler.save("output.pptx")

    也支持从 LLM_Search 结果直接定位:
        filler = MirrorFiller.from_search("DCS 数据中心 架构图")
    """

    def __init__(
        self,
        src_pptx: str | Path,
        slide_num: int = 1,
        layout: str = "left_right",
    ):
        """
        Args:
            src_pptx: 源 PPTX 文件路径
            slide_num: 源 slide 页码 (1-based)
            layout: 布局预设名 ("left_right" / "left_center_right" / ...)
        """
        self.src_pptx = Path(src_pptx)
        self.slide_num = slide_num
        self._layout_name = layout
        self._layout_regions = get_layout(layout) if isinstance(layout, str) else layout
        self._fill_specs: List[Tuple[str, Dict[str, str]]] = []

        # 预加载 slide 信息
        self._prs = Presentation(str(self.src_pptx))
        self._slide = self._prs.slides[slide_num - 1]

        # 缓存 slide XML 名称
        self._slide_xml_name = f"ppt/slides/slide{slide_num}.xml"

    # ── 工厂方法 ──────────────────────────────────

    @classmethod
    def from_pptx(
        cls,
        src_pptx: str | Path,
        slide_num: int = 1,
        layout: str = "left_right",
    ) -> "MirrorFiller":
        """从 PPTX 文件创建"""
        return cls(src_pptx, slide_num, layout)

    @classmethod
    def from_search(
        cls,
        query: str,
        layout: str = "left_right",
        top_k: int = 3,
    ) -> "MirrorFiller":
        """从 LLM_Search 结果自动定位源页。

        需要 LLM_Search 已构建索引。
        """
        import sys
        from pathlib import Path as _Path
        _scripts_dir = _Path(__file__).resolve().parent.parent
        if str(_scripts_dir) not in sys.path:
            sys.path.insert(0, str(_scripts_dir))

        from LLM_Search.search import SemanticSearcher
        searcher = SemanticSearcher()
        results = searcher.search(query, mode="ppt", top_k=top_k)

        if not results.get("ppt"):
            raise ValueError(f'未找到匹配: "{query}"')

        best = results["ppt"][0]
        meta = best["metadata"]
        src_pptx = meta.get("path_abs") or meta.get("file", "")
        slide_num = int(meta.get("slide_num", 1))

        if not Path(src_pptx).exists():
            # 尝试 knowledge_base 路径
            kb_root = _scripts_dir.parent.parent / "knowledge_base"
            for f in kb_root.rglob(Path(src_pptx).name):
                src_pptx = str(f)
                break

        if not Path(src_pptx).exists():
            raise FileNotFoundError(
                f"源文件不存在: {src_pptx}\n"
                f"搜索结果: {meta.get('file')} p{slide_num}"
            )

        print(f'📍 搜索 "{query}" → {Path(src_pptx).name} p{slide_num} (score={best["score"]:.3f})')
        return cls(src_pptx, slide_num, layout)

    # ── 区域探索 ──────────────────────────────────

    def inspect(self) -> None:
        """打印 slide 的区域内容概览，帮助决定填充策略"""
        print(f"\n=== Slide {self.slide_num} 区域概览 ===")
        print(f"布局: {self._layout_name}")
        print(f"shape 总数: {len(self._slide.shapes)}")
        print()

        for region_name, region in self._layout_regions.items():
            matches = _find_text_shapes_in_region(self._slide, region)
            print(f"[{region_name}] ({len(matches)} 文本 shape)")
            for idx, text in matches[:5]:
                print(f"  #{idx}: {text[:120]}")
            if len(matches) > 5:
                print(f"  ... 还有 {len(matches) - 5} 个")
            print()

    def find_text(self, keyword: str) -> List[Tuple[int, str]]:
        """在整个 slide 中搜索包含关键词的 shape"""
        return _find_text_in_slide(self._slide, keyword)

    # ── 填充操作 ──────────────────────────────────

    def fill(
        self,
        layout: str | Dict[str, SlideRegion],
        region_texts: Dict[str, Dict[str, str]],
    ) -> "MirrorFiller":
        """按区域填充文本。

        Args:
            layout: 布局名 ("left_right") 或自定义 Dict[str, SlideRegion]
            region_texts: {区域名: {旧文本: 新文本}}

        Returns:
            self (链式调用)

        Example:
            filler.fill("left_right", {
                "right": {"安全可靠": "企业级稳定性", "极致体验": "统一管理"}
            })
        """
        regions = get_layout(layout) if isinstance(layout, str) else layout

        for region_name, text_map in region_texts.items():
            if region_name not in regions:
                raise ValueError(
                    f'区域 "{region_name}" 不在布局 {list(regions.keys())} 中'
                )
            self._fill_specs.append((region_name, dict(text_map)))

        return self

    # ── 输出 ──────────────────────────────────────

    def save(self, output_path: str | Path) -> Path:
        """执行填充并保存 PPTX。

        工作流:
          1. 复制源 PPTX（保留所有媒体和关系）
          2. 在 ZIP 内删除多余 slide XML，目标 slide 重命名为 slide1
          3. 修正 presentation.xml 的 slide 列表
          4. XML 层执行文本替换
          5. 输出最终文件

        Returns:
            输出文件路径
        """
        import re
        from lxml import etree as _etree

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        tmp_path = output_path.parent / f"_tmp_{output_path.name}"

        # ── Step 1: 复制源 PPTX（完整复制，保证所有资源）──
        shutil.copy2(str(self.src_pptx), str(tmp_path))

        slide_xml_name = f"ppt/slides/slide{self.slide_num}.xml"
        slide_rels_name = f"ppt/slides/_rels/slide{self.slide_num}.xml.rels"

        # ── Step 2: ZIP 层操作 — 删其他 slide + 多余媒体 ──
        with zipfile.ZipFile(tmp_path, 'r') as zin:
            all_files = set(zin.namelist())
            items: dict = {}

            # 2a. 获取目标 slide 引用的媒体文件路径
            kept_media: set = set()
            if slide_rels_name in all_files:
                rels_xml = _etree.fromstring(zin.read(slide_rels_name))
                for rel_el in rels_xml:
                    target = rel_el.get("Target", "")
                    if target and not target.startswith("http"):
                        if target.startswith("../"):
                            kept_media.add("ppt/" + target[3:])
                        else:
                            kept_media.add(f"ppt/slides/{target}")

            # 2b. 过滤文件
            for name in all_files:
                m = re.match(r"ppt/slides/slide(\d+)\.xml$", name)
                if m:
                    if int(m.group(1)) == self.slide_num:
                        items["ppt/slides/slide1.xml"] = zin.read(name)
                    continue

                m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", name)
                if m:
                    if int(m.group(1)) == self.slide_num:
                        items["ppt/slides/_rels/slide1.xml.rels"] = zin.read(name)
                    continue

                # 媒体文件：仅保留目标 slide 引用的
                if "ppt/media/" in name and not name.endswith("/"):
                    if name in kept_media:
                        items[name] = zin.read(name)
                    continue

                # 保留所有其他文件（theme、layout、master、notes 等）
                items[name] = zin.read(name)

        # ── Step 3: 修正 presentation.xml（仅修改 slide 列表）──
        if "ppt/presentation.xml" in items:
            pres_xml = _etree.fromstring(items["ppt/presentation.xml"])
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            sld_id_lst = pres_xml.find(f"{{{ns_p}}}sldIdLst")
            if sld_id_lst is not None:
                kept = None
                for i, child in enumerate(list(sld_id_lst)):
                    if i == self.slide_num - 1:
                        kept = child
                    else:
                        sld_id_lst.remove(child)
                if kept is not None:
                    kept.set("id", "256")
            # 获取目标 slide 的 rId
            target_rId = None
            if kept is not None:
                target_rId = kept.get(f"{{{ns_r}}}id")
            items["ppt/presentation.xml"] = _etree.tostring(
                pres_xml, encoding="UTF-8", xml_declaration=True
            )

        # 修正 presentation.xml.rels：目标 slide 的 Target → slide1.xml
        if target_rId and "ppt/_rels/presentation.xml.rels" in items:
            rels_xml = _etree.fromstring(items["ppt/_rels/presentation.xml.rels"])
            for rel_el in rels_xml:
                if rel_el.get("Id") == target_rId:
                    rel_el.set("Target", "slides/slide1.xml")
                    break
            items["ppt/_rels/presentation.xml.rels"] = _etree.tostring(
                rels_xml, encoding="UTF-8", xml_declaration=True
            )

        # 重建 ZIP
        tmp_path.unlink()
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for name, data in items.items():
                zout.writestr(name, data)

        # ── Step 4: XML 层文本替换 ──
        total_applied = 0
        target_xml_name = "ppt/slides/slide1.xml"
        for region_name, text_map in self._fill_specs:
            regions = (
                get_layout(self._layout_name)
                if isinstance(self._layout_name, str)
                else self._layout_name
            )
            region = regions[region_name]

            region_matches = _find_text_shapes_in_region(self._slide, region)

            region_texts_in_scope = {}
            for old, new in text_map.items():
                for _idx, shape_text in region_matches:
                    if old in shape_text:
                        region_texts_in_scope[old] = new
                        break

            if region_texts_in_scope:
                count = _apply_text_replacements_in_pptx(
                    tmp_path,
                    target_xml_name,
                    region_texts_in_scope,
                )
                total_applied += count
                print(f'  [{region_name}] 替换 {count} 处文本')
            else:
                print(f'  [{region_name}] ⚠️ 区域内未找到匹配文本')

        # 清理并输出
        if output_path.exists():
            output_path.unlink()
        shutil.move(str(tmp_path), str(output_path))

        tmp_files = list(output_path.parent.glob(f"_tmp_*"))
        for f in tmp_files:
            f.unlink(missing_ok=True)

        print(f"\n✅ 已保存: {output_path} (共替换 {total_applied} 处)")
        return output_path

    # ── 信息 ──────────────────────────────────────

    def info(self) -> dict:
        """返回 filler 状态信息"""
        return {
            "src_pptx": str(self.src_pptx),
            "slide_num": self.slide_num,
            "layout": self._layout_name,
            "fill_specs": [
                {"region": r, "count": len(t)}
                for r, t in self._fill_specs
            ],
        }
