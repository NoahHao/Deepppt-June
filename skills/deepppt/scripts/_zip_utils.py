#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""公共工具：单页/多页 zip-filter 提取 + 清空→填充。修复 rels 过滤 bug。"""
import os, re, zipfile
from xml.etree import ElementTree as ET

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

for prefix, uri in [('p', NS_P), ('r', NS_R)]:
    ET.register_namespace(prefix, uri)
ET.register_namespace('', NS_P)


def extract_slides(template_path, output_path, slide_numbers):
    """
    从模板提取指定页为独立PPTX。
    输出 PPTX 中的 slide 始终重命名为 slide1.xml（带 slide1.xml.rels）。

    修复：
    - presentation.xml.rels 只删除 slide 关系，保留 slideMaster/theme/notes
    - 输出 slide 始终为 slide1.xml，便于下游合并
    """
    slide_numbers = list(slide_numbers)

    with zipfile.ZipFile(template_path, 'r') as zin:
        # 读取 slide rId 映射
        rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        root_rel = ET.fromstring(rels_xml)
        slide_rid = {}
        slide_targets = {}  # rId -> target filename
        for child in root_rel:
            target = child.get('Target', '')
            m = re.search(r'slide(\d+)', target)
            if m and 'slide' in child.get('Type', ''):
                slide_rid[int(m.group(1))] = child.get('Id')
                slide_targets[child.get('Id')] = target

        keep_rids = {slide_rid[sn] for sn in slide_numbers if sn in slide_rid}

        # Keep slides: original numbers
        keep_slide_nums = set(slide_numbers)

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                fn = item.filename
                data = zin.read(fn)

                # Skip non-kept slides
                m = re.match(r'ppt/slides/slide(\d+)\.xml$', fn)
                if m and int(m.group(1)) not in keep_slide_nums:
                    continue
                m2 = re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels$', fn)
                if m2 and int(m2.group(1)) not in keep_slide_nums:
                    continue

                # ── Rename kept slide to slide1.xml ──
                if m and int(m.group(1)) in keep_slide_nums:
                    old_num = int(m.group(1))
                    if old_num != 1:
                        # Write as slide1.xml
                        zout.writestr(zipfile.ZipInfo('ppt/slides/slide1.xml'), data)
                        continue

                if m2 and int(m2.group(1)) in keep_slide_nums:
                    old_num = int(m2.group(1))
                    if old_num != 1:
                        # Write as slide1.xml.rels
                        zout.writestr(zipfile.ZipInfo('ppt/slides/_rels/slide1.xml.rels'), data)
                        continue

                # ── Fix presentation.xml: rename rId to rId1, target to slide1.xml ──
                if fn == 'ppt/presentation.xml':
                    root_p = ET.fromstring(data)
                    sld_lst = root_p.find('.//{' + NS_P + '}sldIdLst')
                    if sld_lst is not None:
                        for c in list(sld_lst):
                            sld_lst.remove(c)
                        # Add single slide referencing rId1
                        sld_id = ET.SubElement(sld_lst, 'p:sldId')
                        sld_id.set('id', '256')
                        sld_id.set('r:id', 'rId1')
                    data = ET.tostring(root_p, encoding='UTF-8', xml_declaration=True)

                # ── Fix presentation.xml.rels: keep only master/layout/theme, add single slide rel rId1→slide1.xml ──
                elif fn == 'ppt/_rels/presentation.xml.rels':
                    root_rel2 = ET.fromstring(data)
                    keep_children = []
                    for c in list(root_rel2):
                        typ = c.get('Type', '')
                        # Remove ALL slide relationships (we'll add one for slide1)
                        if 'slide' in typ and 'slides/slide' in c.get('Target', ''):
                            continue
                        keep_children.append(c)
                    
                    # Build new rels root
                    ns_rel = 'http://schemas.openxmlformats.org/package/2006/relationships'
                    new_root = ET.Element(f'{{{ns_rel}}}Relationships')
                    for c in keep_children:
                        new_root.append(c)
                    
                    # Add single slide relationship → slide1.xml
                    slide_rel = ET.SubElement(new_root, 'Relationship')
                    slide_rel.set('Id', 'rId1')
                    slide_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
                    slide_rel.set('Target', 'slides/slide1.xml')
                    
                    data = ET.tostring(new_root, encoding='UTF-8', xml_declaration=True)

                # ── Fix Content_Types: remove old slide entries, add slide1 ──
                elif fn == '[Content_Types].xml':
                    root_ct = ET.fromstring(data)
                    for child in list(root_ct):
                        pn = child.get('PartName', '')
                        m = re.search(r'/ppt/slides/slide(\d+)\.xml', pn)
                        if m:
                            if int(m.group(1)) not in keep_slide_nums:
                                root_ct.remove(child)
                            elif int(m.group(1)) != 1:
                                # Rename to /ppt/slides/slide1.xml
                                child.set('PartName', '/ppt/slides/slide1.xml')
                    data = ET.tostring(root_ct, encoding='UTF-8', xml_declaration=True)

                zout.writestr(item, data)

    print(f"   提取 {len(slide_numbers)} 页: {slide_numbers}")



def _slot_capacity(shape, para):
    """估算单个段落的容量（可容纳字数）"""
    from pptx.util import Emu, Inches
    w_inches = (shape.width or Emu(0)) / 914400.0
    h_inches = (shape.height or Emu(0)) / 914400.0
    fs = 12.0
    if para.runs:
        rpr = para.runs[0].font
        fs = rpr.size / 12700.0 if rpr.size else 12.0

    if fs <= 0 or w_inches <= 0:
        return 0, 0

    chars_per_line = max(1, int(w_inches * 72 / (fs * 0.55)))
    line_height_inches = fs * 1.4 / 72
    lines = max(1, int(h_inches / line_height_inches)) if h_inches > 0 else 1
    capacity = chars_per_line * lines
    return capacity, chars_per_line


def fill_slide_sequential(slide, new_content):
    """
    清空所有文本 → 按位置顺序填入 new_content。
    ⚠️ 简单顺序填充，不检查容量。推荐使用 fill_slide_smart()。
    返回 (填充数, 总槽数, 原文列表)。
    """
    from pptx.util import Emu

    items = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    items.append((shape, para, pi, t))

    items.sort(key=lambda x: (x[0].top or Emu(0), x[0].left or Emu(0), x[2]))

    for _, para, _, _ in items:
        for run in para.runs:
            run.text = ""

    filled = 0
    for i, (_, para, _, orig) in enumerate(items):
        if i < len(new_content):
            para.runs[0].text = new_content[i]
            filled += 1
            print(f"   [{i}] '{orig[:35]}' → '{new_content[i][:35]}'")

    return filled, len(items), items


def fill_slide_smart(slide, new_content, verbose=True):
    """
    清空所有文本 → 按槽位容量智能匹配填入 new_content。

    算法：
      1. 计算每槽容量（字号×宽度×行数）
      2. 槽位按容量降序排列，内容按字数降序排列
      3. 最长内容匹配到最大容量槽位
      4. 内容超容量时 ⚠️ 告警

    返回 (filled, total, warnings[]).
    """
    from pptx.util import Emu

    items = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    cap, cpl = _slot_capacity(shape, para)
                    items.append({
                        'shape': shape, 'para': para, 'pi': pi,
                        'orig': t, 'capacity': cap, 'cpl': cpl,
                        'top': (shape.top or Emu(0)),
                        'left': (shape.left or Emu(0)),
                    })

    # 按容量降序 → 大槽在前
    items.sort(key=lambda x: (-x['capacity'], x['top'], x['left']))
    # 内容按字数降序 → 长文在前
    indexed = [(len(c), i, c) for i, c in enumerate(new_content)]
    indexed.sort(key=lambda x: -x[0])

    # 清空全部
    for it in items:
        for run in it['para'].runs:
            run.text = ""

    # 打印槽位信息
    if verbose:
        print(f"\n  {'槽位容量':-^60}")
        for i, it in enumerate(items):
            flag = "⚠ 超" if (indexed[i][0] if i < len(indexed) else 0) > it['capacity'] else "  "
            print(f"  [{i}] 容{it['capacity']:>3}字 | {it['orig'][:25]:<25} | → ", end="")

    # 智能匹配：最长→最大槽
    warnings = []
    filled = 0
    for slot_idx, it in enumerate(items):
        if slot_idx < len(indexed):
            _, content_idx, text = indexed[slot_idx]
            it['para'].runs[0].text = text
            filled += 1
            if verbose:
                flag = "⚠超" if len(text) > it['capacity'] else "  ✓"
                print(f"{flag} [{len(text):>2}字] {text[:40]}")
            if len(text) > it['capacity']:
                warnings.append(
                    f"  ⚠ 槽[{slot_idx}] 容量{it['capacity']}字 ← 填入{len(text)}字 '{text[:30]}'"
                )
        else:
            if verbose:
                print(f"     (清空)")

    if warnings:
        print(f"\n  ⚠️ 容量超限 {len(warnings)} 处:")
        for w in warnings:
            print(w)
    else:
        print(f"\n  ✅ 全部在容量范围内")

    return filled, len(items), warnings


def verify_output(pptx_path, expected_content):
    """
    Bug自查：验证输出PPTX文本无残文。
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    issues = []
    for si, slide in enumerate(prs.slides):
        actual = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        actual.append(t)

        # 检查残文（实际文本不在预期内容中）
        for t in actual:
            if t not in expected_content:
                issues.append(f"  页{si+1} 残文: '{t[:60]}'")

    if issues:
        print("\n⚠️ Bug自查 FAILED:")
        for i in issues:
            print(i)
    else:
        print(f"\n✅ Bug自查 PASS: {len(prs.slides)}页, 0残文")
    return len(issues) == 0
