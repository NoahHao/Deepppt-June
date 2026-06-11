#!/usr/bin/env python3
"""生成一页介绍小狗的 PPT 幻灯片"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── 创建 PPT ──
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# ── 背景 ──
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)  # 温暖米白

# ── 左侧装饰色块 ──
from pptx.util import Inches, Pt
left_block = slide.shapes.add_shape(
    1,  # MSO_SHAPE.RECTANGLE
    Inches(0), Inches(0),
    Inches(1.2), Inches(7.5),
)
left_block.fill.solid()
left_block.fill.fore_color.rgb = RGBColor(0xFF, 0x9A, 0x76)  # 珊瑚橙
left_block.line.fill.background()

# ── 标题 ──
title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.5), Inches(9), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "认识人类最好的朋友 — 小狗"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
p.font.name = "Microsoft YaHei"

# 副标题
sub_box = slide.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9), Inches(0.6))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.text = "🐶 犬科动物 · 人类驯化已有 15,000+ 年历史"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
p2.font.name = "Microsoft YaHei"

# ── 图片（右侧）──
img_path = os.path.join(os.path.dirname(__file__), "cute_dog.png")
if os.path.exists(img_path):
    pic = slide.shapes.add_picture(
        img_path,
        Inches(8.3), Inches(2.2),
        Inches(4.5), Inches(4.5),
    )

# ── 左侧文字内容 ──
content_data = [
    ("🐾 品种多样性", "全球超过 340 种犬类品种\n从茶杯犬到藏獒，体型差异可达 100 倍"),
    ("❤️ 忠诚与陪伴", "狗狗是人类最忠诚的伙伴\n能感知主人情绪，缓解压力与孤独"),
    ("🧠 超凡智慧", "最聪明的犬种可理解 250+ 个单词\n相当于 2 岁幼儿的认知水平"),
    ("⚕️ 健康守护者", "导盲犬、搜救犬、治疗犬\n在医疗和公共服务中发挥重要作用"),
]

y_start = 2.2
for i, (title, desc) in enumerate(content_data):
    y = y_start + i * 1.2

    # 图标 + 标题
    item_box = slide.shapes.add_textbox(Inches(1.8), Inches(y), Inches(5.5), Inches(1.0))
    tf_item = item_box.text_frame
    tf_item.word_wrap = True

    p_title = tf_item.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    p_title.font.name = "Microsoft YaHei"

    p_desc = tf_item.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    p_desc.font.name = "Microsoft YaHei"
    p_desc.space_before = Pt(2)

# ── 底部分隔线 ──
line = slide.shapes.add_shape(
    1, Inches(1.8), Inches(6.6), Inches(7), Pt(2),
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xE8, 0xE0, 0xD8)
line.line.fill.background()

# ── 页脚 ──
footer_box = slide.shapes.add_textbox(Inches(1.8), Inches(6.7), Inches(9), Inches(0.5))
tf_f = footer_box.text_frame
p_f = tf_f.paragraphs[0]
p_f.text = "数据来源：AKC / FCI · 图片由 AI 生成"
p_f.font.size = Pt(10)
p_f.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p_f.font.name = "Microsoft YaHei"
p_f.alignment = PP_ALIGN.LEFT

# ── 保存 ──
output_path = os.path.join(os.path.dirname(__file__), "小狗介绍.pptx")
prs.save(output_path)
print(f"[SAVED] {output_path}")
